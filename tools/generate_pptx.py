"""
generate_pptx.py

Generates a PowerPoint presentation explaining the WAT Workflow Agent,
including graphical sequence diagrams drawn with python-pptx shapes.

Output: docs/WAT_Workflow_Agent.pptx

Usage:
    python tools/generate_pptx.py
"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


# ── Colour palette ────────────────────────────────────────────────────────────
C_NAVY    = RGBColor(0x1a, 0x2a, 0x4a)
C_BLUE    = RGBColor(0x29, 0x80, 0xb9)
C_GREEN   = RGBColor(0x27, 0xae, 0x60)
C_ORANGE  = RGBColor(0xf3, 0x9c, 0x12)
C_RED     = RGBColor(0xc0, 0x39, 0x2b)
C_TEAL    = RGBColor(0x00, 0xbc, 0xd4)
C_WHITE   = RGBColor(0xff, 0xff, 0xff)
C_LGREY   = RGBColor(0xf0, 0xf4, 0xf8)
C_DGREY   = RGBColor(0x55, 0x55, 0x55)
C_ARROW   = RGBColor(0x33, 0x33, 0x33)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


# ── Helpers ───────────────────────────────────────────────────────────────────

def add_slide(prs, layout_idx=6):
    layout = prs.slide_layouts[layout_idx]
    return prs.slides.add_slide(layout)


def rect(slide, x, y, w, h, fill, text="", font_size=11,
         font_color=C_WHITE, bold=False, align=PP_ALIGN.CENTER,
         radius=False, border_color=None, border_width=Pt(0)):
    from pptx.util import Pt
    from pptx.oxml.ns import qn
    from lxml import etree

    shape = slide.shapes.add_shape(
        1 if not radius else 5,  # MSO_SHAPE_TYPE: rectangle=1, rounded=5
        x, y, w, h
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = border_width
    else:
        shape.line.fill.background()

    if text:
        tf = shape.text_frame
        tf.word_wrap = True
        tf.auto_size = None
        from pptx.enum.text import MSO_ANCHOR
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.size = Pt(font_size)
        run.font.color.rgb = font_color
        run.font.bold = bold
    return shape


def label(slide, x, y, w, h, text, font_size=11,
          font_color=C_DGREY, bold=False, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.color.rgb = font_color
    run.font.bold = bold
    return tb


def arrow(slide, x1, y1, x2, y2, color=C_ARROW, width=Pt(1.5), label_text="", dashed=False):
    """Draw a horizontal or angled arrow from (x1,y1) to (x2,y2)."""
    from pptx.util import Pt
    from pptx.oxml.ns import qn
    connector = slide.shapes.add_connector(1, x1, y1, x2, y2)  # 1=straight
    connector.line.color.rgb = color
    connector.line.width = width
    if dashed:
        connector.line.dash_style = 4  # dash
    # arrowhead
    from lxml import etree
    ln = connector.line._ln
    tail = ln.find(qn("a:tailEnd"))
    if tail is None:
        tail = etree.SubElement(ln, qn("a:tailEnd"))
    tail.set("type", "none")
    head = ln.find(qn("a:headEnd"))
    if head is None:
        head = etree.SubElement(ln, qn("a:headEnd"))
    head.set("type", "arrow")
    head.set("w", "med")
    head.set("len", "med")

    if label_text:
        mid_x = (x1 + x2) / 2 - Inches(0.6)
        mid_y = min(y1, y2) - Inches(0.18)
        label(slide, mid_x, mid_y, Inches(1.2), Inches(0.25),
              label_text, font_size=8, font_color=C_DGREY)


def lifeline(slide, x, y_top, y_bot, color=C_DGREY):
    """Draw a vertical dashed lifeline."""
    from pptx.util import Pt
    line = slide.shapes.add_connector(1, x, y_top, x, y_bot)
    line.line.color.rgb = color
    line.line.width = Pt(1)
    line.line.dash_style = 4


# ── Slide builders ────────────────────────────────────────────────────────────

def slide_title(prs):
    slide = add_slide(prs)
    # Background
    rect(slide, 0, 0, SLIDE_W, SLIDE_H, C_NAVY)
    # Accent bar
    rect(slide, 0, Inches(4.5), SLIDE_W, Inches(0.08), C_BLUE)
    # Title
    label(slide, Inches(1), Inches(1.5), Inches(11), Inches(1.5),
          "WAT Workflow Agent", font_size=52, font_color=C_WHITE, bold=True,
          align=PP_ALIGN.CENTER)
    # Subtitle
    label(slide, Inches(1), Inches(3.2), Inches(11), Inches(0.8),
          "Workflows · Agents · Tools",
          font_size=22, font_color=RGBColor(0xa0, 0xc8, 0xf0),
          align=PP_ALIGN.CENTER)
    label(slide, Inches(1), Inches(4.8), Inches(11), Inches(0.6),
          "Automate business process diagrams from plain English — via Slack or email",
          font_size=14, font_color=RGBColor(0xcc, 0xdd, 0xee),
          align=PP_ALIGN.CENTER)


def slide_what_it_does(prs):
    slide = add_slide(prs)
    rect(slide, 0, 0, SLIDE_W, SLIDE_H, C_LGREY)
    rect(slide, 0, 0, SLIDE_W, Inches(1.1), C_NAVY)
    label(slide, Inches(0.4), Inches(0.2), Inches(12), Inches(0.7),
          "What the WAT Workflow Agent Does",
          font_size=28, font_color=C_WHITE, bold=True)

    items = [
        (C_BLUE,   "📝  Accepts plain-English process descriptions",
                   "Describe any business process in natural language — no flowchart tools needed."),
        (C_GREEN,  "🤖  Builds Miro workflow diagrams automatically",
                   "Creates a professional swim-lane diagram on a Miro board with nodes, connectors, and colour-coded lanes."),
        (C_ORANGE, "📸  Screenshots and delivers the result",
                   "Takes a Playwright screenshot of the board and sends it via Slack thread or email."),
        (C_TEAL,   "✏️  Accepts change requests",
                   "Say 'add a step', 'rename X', or 'remove Y' and the agent updates the live Miro board."),
    ]

    for i, (color, heading, body) in enumerate(items):
        col = i % 2
        row = i // 2
        bx = Inches(0.4 + col * 6.4)
        by = Inches(1.4 + row * 2.7)
        rect(slide, bx, by, Inches(6.0), Inches(2.4), C_WHITE,
             border_color=color, border_width=Pt(2))
        rect(slide, bx, by, Inches(6.0), Inches(0.5), color)
        label(slide, bx + Inches(0.15), by + Inches(0.05), Inches(5.7), Inches(0.45),
              heading, font_size=13, font_color=C_WHITE, bold=True)
        label(slide, bx + Inches(0.2), by + Inches(0.6), Inches(5.6), Inches(1.6),
              body, font_size=11, font_color=C_DGREY)


def slide_architecture(prs):
    slide = add_slide(prs)
    rect(slide, 0, 0, SLIDE_W, SLIDE_H, C_LGREY)
    rect(slide, 0, 0, SLIDE_W, Inches(1.1), C_NAVY)
    label(slide, Inches(0.4), Inches(0.2), Inches(12), Inches(0.7),
          "The WAT Architecture",
          font_size=28, font_color=C_WHITE, bold=True)

    layers = [
        (C_BLUE,   "Layer 1: Workflows",   "Markdown SOPs in workflows/\n\nDefine objectives, steps, edge cases, and tool sequences in plain language."),
        (C_GREEN,  "Layer 2: Agent",        "Claude (this agent)\n\nReads workflows, makes decisions, calls tools in sequence, recovers from errors."),
        (C_ORANGE, "Layer 3: Tools",        "Python scripts in tools/\n\nDeterministic execution: API calls, file I/O, data transforms. Credentials in .env."),
    ]

    arrow_y = Inches(4.0)
    for i, (color, title, body) in enumerate(layers):
        bx = Inches(0.5 + i * 4.2)
        by = Inches(1.4)
        rect(slide, bx, by, Inches(3.8), Inches(4.8), C_WHITE,
             border_color=color, border_width=Pt(2))
        rect(slide, bx, by, Inches(3.8), Inches(0.65), color)
        label(slide, bx + Inches(0.15), by + Inches(0.1), Inches(3.5), Inches(0.5),
              title, font_size=14, font_color=C_WHITE, bold=True)
        label(slide, bx + Inches(0.2), by + Inches(0.85), Inches(3.4), Inches(3.7),
              body, font_size=11, font_color=C_DGREY)
        if i < 2:
            arrow(slide,
                  bx + Inches(3.8), arrow_y,
                  bx + Inches(4.2), arrow_y,
                  color=C_BLUE, width=Pt(2.5))

    label(slide, Inches(0.5), Inches(6.5), Inches(12), Inches(0.5),
          "Why separate layers?  Each tool call is 100% reliable; only the agent's reasoning is probabilistic.",
          font_size=10, font_color=C_DGREY, align=PP_ALIGN.CENTER)


def _seq_actors(slide, actors, y_actor, actor_w, actor_h, lifeline_bot):
    """Draw actor boxes and lifelines. Returns list of center-x per actor."""
    centers = []
    total_w = SLIDE_W - Inches(0.6)
    step = total_w / len(actors)
    for i, (name, color) in enumerate(actors):
        cx = Inches(0.3) + step * i + step / 2
        bx = cx - actor_w / 2
        rect(slide, bx, y_actor, actor_w, actor_h, color,
             text=name, font_size=9, bold=True, font_color=C_WHITE)
        lifeline(slide, cx, y_actor + actor_h, lifeline_bot)
        centers.append(cx)
    return centers


def _seq_msg(slide, centers, from_i, to_i, y, text, ret=False, color=C_ARROW):
    cx_from = centers[from_i]
    cx_to   = centers[to_i]
    arrow(slide, cx_from, y, cx_to, y,
          color=color, width=Pt(1.5 if not ret else 1.0),
          dashed=ret)
    # label above the arrow
    mid_x = (min(cx_from, cx_to) + abs(cx_to - cx_from) / 2) - Inches(0.8)
    label(slide, mid_x, y - Inches(0.2), Inches(1.6), Inches(0.22),
          text, font_size=7.5, font_color=C_DGREY, align=PP_ALIGN.CENTER)


def slide_seq_slack(prs):
    slide = add_slide(prs)
    rect(slide, 0, 0, SLIDE_W, SLIDE_H, C_LGREY)
    rect(slide, 0, 0, SLIDE_W, Inches(0.9), C_NAVY)
    label(slide, Inches(0.4), Inches(0.12), Inches(12), Inches(0.65),
          "Sequence: Slack Trigger Flow",
          font_size=24, font_color=C_WHITE, bold=True)

    actor_y   = Inches(1.0)
    actor_w   = Inches(1.3)
    actor_h   = Inches(0.5)
    life_bot  = Inches(7.2)

    actors = [
        ("User",        C_NAVY),
        ("Slack",       RGBColor(0x4a, 0x15, 0x4b)),
        ("Listener",    C_BLUE),
        ("Thread",      C_TEAL),
        ("Parse",       C_GREEN),
        ("MiroCreate",  C_ORANGE),
        ("MiroRender",  RGBColor(0xe6, 0x7e, 0x22)),
        ("Screenshot",  C_RED),
        ("Respond",     RGBColor(0x16, 0xa0, 0x85)),
    ]
    centers = _seq_actors(slide, actors, actor_y, actor_w, actor_h, life_bot)

    msgs = [
        (0, 1, "@mention: process desc",  False),
        (1, 2, "app_mention event",        False),
        (2, 1, "acknowledge (<3s)",         True),
        (2, 1, "post 'Building...'",        False),
        (2, 3, "spawn daemon thread",       False),
        (3, 4, "parse_workflow.py",         False),
        (4, 3, "workflow.json",             True),
        (3, 5, "miro_create_board.py",      False),
        (5, 3, "board_info.json",           True),
        (3, 6, "miro_add_swimlane.py",      False),
        (6, 3, "shapes + connectors",       True),
        (3, 7, "screenshot_board.py",       False),
        (7, 3, "workflow.png",              True),
        (3, 8, "slack_respond.py",          False),
        (8, 1, "board link + PNG in thread",False),
    ]

    y_start = Inches(1.75)
    row_h   = Inches(0.34)
    for i, (fi, ti, txt, ret) in enumerate(msgs):
        _seq_msg(slide, centers, fi, ti, y_start + row_h * i, txt, ret=ret)


def slide_seq_update(prs):
    slide = add_slide(prs)
    rect(slide, 0, 0, SLIDE_W, SLIDE_H, C_LGREY)
    rect(slide, 0, 0, SLIDE_W, Inches(0.9), C_NAVY)
    label(slide, Inches(0.4), Inches(0.12), Inches(12), Inches(0.65),
          "Sequence: Change Request & Manual/Email Flow",
          font_size=24, font_color=C_WHITE, bold=True)

    # --- Update flow (left half) ---
    rect(slide, Inches(0.3), Inches(1.0), Inches(6.1), Inches(6.2),
         C_WHITE, border_color=C_BLUE, border_width=Pt(1))
    label(slide, Inches(0.5), Inches(1.05), Inches(5), Inches(0.35),
          "Change Request Flow", font_size=12, font_color=C_BLUE, bold=True)

    upd_actors = [
        ("User",    C_NAVY),
        ("Slack",   RGBColor(0x4a, 0x15, 0x4b)),
        ("Listener",C_BLUE),
        ("Update",  C_ORANGE),
        ("Miro",    C_TEAL),
        ("Respond", RGBColor(0x16, 0xa0, 0x85)),
    ]
    actor_y  = Inches(1.5)
    actor_w  = Inches(0.88)
    actor_h  = Inches(0.42)
    life_bot = Inches(7.1)

    total_w = Inches(5.8)
    step = total_w / len(upd_actors)
    u_centers = []
    for i, (name, color) in enumerate(upd_actors):
        cx = Inches(0.4) + step * i + step / 2
        bx = cx - actor_w / 2
        rect(slide, bx, actor_y, actor_w, actor_h, color,
             text=name, font_size=8, bold=True, font_color=C_WHITE)
        lifeline(slide, cx, actor_y + actor_h, life_bot)
        u_centers.append(cx)

    upd_msgs = [
        (0, 1, "add/change/remove...", False),
        (1, 2, "app_mention",          False),
        (2, 3, "miro_update.py",       False),
        (3, 4, "Claude: interpret",    False),
        (4, 3, "edit shapes",          True),
        (3, 5, "slack_respond.py",     False),
        (5, 1, "updated PNG + link",   False),
    ]
    y_start = Inches(2.1)
    row_h   = Inches(0.6)
    for i, (fi, ti, txt, ret) in enumerate(upd_msgs):
        _seq_msg(slide, u_centers, fi, ti, y_start + row_h * i, txt, ret=ret)

    # --- Manual/Email flow (right half) ---
    rect(slide, Inches(6.8), Inches(1.0), Inches(6.2), Inches(6.2),
         C_WHITE, border_color=C_GREEN, border_width=Pt(1))
    label(slide, Inches(7.0), Inches(1.05), Inches(5), Inches(0.35),
          "Manual / Email Flow", font_size=12, font_color=C_GREEN, bold=True)

    em_actors = [
        ("User",   C_NAVY),
        ("Agent",  C_BLUE),
        ("Parse",  C_GREEN),
        ("Miro",   C_ORANGE),
        ("Shot",   C_RED),
        ("Resend", RGBColor(0x16, 0xa0, 0x85)),
    ]
    total_w2 = Inches(5.8)
    step2 = total_w2 / len(em_actors)
    e_centers = []
    for i, (name, color) in enumerate(em_actors):
        cx = Inches(7.0) + step2 * i + step2 / 2
        bx = cx - actor_w / 2
        rect(slide, bx, actor_y, actor_w, actor_h, color,
             text=name, font_size=8, bold=True, font_color=C_WHITE)
        lifeline(slide, cx, actor_y + actor_h, life_bot)
        e_centers.append(cx)

    em_msgs = [
        (0, 1, "describe process",     False),
        (1, 2, "parse_workflow.py",    False),
        (2, 1, "workflow.json",        True),
        (1, 3, "create + render",      False),
        (3, 1, "board_url",            True),
        (1, 4, "screenshot_board.py",  False),
        (4, 1, "workflow.png",         True),
        (1, 5, "send_email.py",        False),
        (5, 0, "email + PNG",          False),
    ]
    y_start2 = Inches(2.1)
    for i, (fi, ti, txt, ret) in enumerate(em_msgs):
        _seq_msg(slide, e_centers, fi, ti, y_start2 + row_h * i, txt, ret=ret)


def slide_tools(prs):
    slide = add_slide(prs)
    rect(slide, 0, 0, SLIDE_W, SLIDE_H, C_LGREY)
    rect(slide, 0, 0, SLIDE_W, Inches(1.1), C_NAVY)
    label(slide, Inches(0.4), Inches(0.2), Inches(12), Inches(0.7),
          "Key Tools & Files",
          font_size=28, font_color=C_WHITE, bold=True)

    tools = [
        (C_GREEN,  "parse_swimlane_workflow.py",  "Claude API → swim lane workflow JSON with lane/col assignments"),
        (C_BLUE,   "miro_create_board.py",         "Miro API → creates a new board, saves board_id + URL"),
        (C_BLUE,   "miro_add_swimlane.py",          "Miro API → renders shapes, connectors, and lane backgrounds"),
        (C_ORANGE, "miro_update_workflow.py",       "Claude + Miro API → applies targeted node/edge edits"),
        (C_RED,    "screenshot_board.py",           "Playwright headless Chromium → login, screenshot, Pillow crop"),
        (C_TEAL,   "slack_listener.py",             "Socket Mode listener → routes mentions to create/update/rerender"),
        (C_TEAL,   "slack_respond.py",              "Posts board link + uploads PNG and .mmd file to Slack thread"),
        (C_DGREY,  "generate_mermaid.py",           "Converts workflow JSON → Mermaid flowchart (.mmd) for VS Code/GitHub"),
        (RGBColor(0x8e,0x44,0xad), "send_email.py","Resend API → transactional email with base64 PNG attachment"),
    ]

    col_w = Inches(6.1)
    for i, (color, name, desc) in enumerate(tools):
        col = i % 2
        row = i // 2
        bx = Inches(0.35) + col * col_w
        by = Inches(1.3) + row * Inches(1.45)
        rect(slide, bx, by, col_w - Inches(0.1), Inches(1.3),
             C_WHITE, border_color=color, border_width=Pt(2))
        rect(slide, bx, by, col_w - Inches(0.1), Inches(0.4), color)
        label(slide, bx + Inches(0.12), by + Inches(0.05),
              col_w - Inches(0.3), Inches(0.35),
              name, font_size=10, font_color=C_WHITE, bold=True)
        label(slide, bx + Inches(0.15), by + Inches(0.5),
              col_w - Inches(0.3), Inches(0.75),
              desc, font_size=10, font_color=C_DGREY)


def slide_summary(prs):
    slide = add_slide(prs)
    rect(slide, 0, 0, SLIDE_W, SLIDE_H, C_NAVY)
    rect(slide, 0, Inches(2.8), SLIDE_W, Inches(0.06), C_BLUE)

    label(slide, Inches(1), Inches(0.6), Inches(11), Inches(1.0),
          "Ready to Use", font_size=40, font_color=C_WHITE, bold=True,
          align=PP_ALIGN.CENTER)

    steps = [
        ("1", C_BLUE,   "Create Slack app →\nadd .env tokens"),
        ("2", C_GREEN,  "python tools/\nslack_listener.py"),
        ("3", C_ORANGE, "@mention the bot\nwith any process"),
        ("4", C_TEAL,   "Receive Miro board\n+ PNG + .mmd in thread"),
    ]
    for i, (num, color, text) in enumerate(steps):
        bx = Inches(0.8 + i * 3.0)
        by = Inches(3.3)
        rect(slide, bx, by, Inches(2.5), Inches(2.8), color,
             border_color=C_WHITE, border_width=Pt(1))
        label(slide, bx + Inches(0.1), by + Inches(0.15),
              Inches(2.3), Inches(0.6),
              num, font_size=32, font_color=C_WHITE, bold=True,
              align=PP_ALIGN.CENTER)
        label(slide, bx + Inches(0.1), by + Inches(0.9),
              Inches(2.3), Inches(1.6),
              text, font_size=13, font_color=C_WHITE,
              align=PP_ALIGN.CENTER)

    label(slide, Inches(1), Inches(6.5), Inches(11), Inches(0.6),
          "Works locally · No server required · Socket Mode WebSocket connection",
          font_size=12, font_color=RGBColor(0xa0, 0xc0, 0xe0),
          align=PP_ALIGN.CENTER)


# ── Main ──────────────────────────────────────────────────────────────────────

def main():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    slide_title(prs)
    slide_what_it_does(prs)
    slide_architecture(prs)
    slide_seq_slack(prs)
    slide_seq_update(prs)
    slide_tools(prs)
    slide_summary(prs)

    out = Path("docs/WAT_Workflow_Agent.pptx")
    out.parent.mkdir(exist_ok=True)
    prs.save(str(out))
    print(f"Saved: {out}")


if __name__ == "__main__":
    main()
